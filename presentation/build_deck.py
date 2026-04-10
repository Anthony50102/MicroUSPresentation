"""
Full 20-slide deck v4 — Gradient Noir
"The Next Layer for Digital Twins: From Simulation to Autonomous Action"

Incorporates all feedback from review round.
v4.1: Build animations (S5, S9), teal accent (S11, S18), native arrow
connectors, speaker notes on all slides, high-res QR codes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import math
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Gradient Noir Palette ──
BG1 = '#0F172A'
BG2 = '#1E1E2E'
WHITE = '#F5F5F0'
CREAM = '#E8E4D9'
DIM = '#8B919C'
DARKER_DIM = '#6B7280'
ACCENT = '#D4D0C8'
TEAL = '#2DD4BF'
CARD_BG = '#151D30'
CARD_BG2 = '#1A2236'
FONT = 'Segoe UI'


def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = hex_to_rgb(BG1)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = hex_to_rgb(BG2)
    fill.gradient_stops[1].position = 1.0


def tb(slide, left, top, width, height, text, size, color,
       bold=False, align=PP_ALIGN.LEFT, italic=False, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = hex_to_rgb(color)
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return txBox


def multi_para(slide, left, top, width, height, lines, font=FONT,
               align=PP_ALIGN.LEFT, spacing=6):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = hex_to_rgb(color)
        p.font.bold = bold
        p.alignment = align
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
    return txBox


def rect(slide, left, top, width, height, fill=None, line=None, lw=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = hex_to_rgb(line)
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def rrect(slide, left, top, width, height, fill=None, line=None, lw=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = hex_to_rgb(line)
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def rrect_text(slide, left, top, width, height, text, size, color,
               fill=None, line=None, bold=True, align=PP_ALIGN.CENTER):
    shape = rrect(slide, left, top, width, height, fill, line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.color.rgb = hex_to_rgb(color)
    p.font.bold = bold
    p.alignment = align
    return shape


def oval(slide, left, top, w, h, fill=None, line=None, lw=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = hex_to_rgb(line)
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def oval_text(slide, cx, cy, r, text, size, font_color,
              fill=None, line=None, bold=True):
    shape = oval(slide, cx - r, cy - r, r * 2, r * 2, fill, line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return shape


def connector(slide, x1, y1, x2, y2, color, width=Pt(1.5)):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = hex_to_rgb(color)
    c.line.width = width
    return c


def arrow_connector(slide, x1, y1, x2, y2, color, width=Pt(1.5)):
    """Straight connector with arrowhead at the end."""
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = hex_to_rgb(color)
    c.line.width = width
    sp_pr = c._element.find(qn('p:spPr'))
    ln = sp_pr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(sp_pr, qn('a:ln'))
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return c


def add_appear_animations(slide, shape_groups):
    """Add click-to-appear build animations to a slide.
    shape_groups = [[shapes_click_1], [shapes_click_2], ...]
    Each inner list appears together on one click.
    """
    timing = etree.SubElement(slide._element, qn('p:timing'))
    tnLst = etree.SubElement(timing, qn('p:tnLst'))
    root_par = etree.SubElement(tnLst, qn('p:par'))
    root_cTn = etree.SubElement(root_par, qn('p:cTn'))
    root_cTn.set('id', '1')
    root_cTn.set('dur', 'indefinite')
    root_cTn.set('restart', 'never')
    root_cTn.set('nodeType', 'tmRoot')
    root_children = etree.SubElement(root_cTn, qn('p:childTnLst'))

    seq = etree.SubElement(root_children, qn('p:seq'))
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')
    seq_cTn = etree.SubElement(seq, qn('p:cTn'))
    seq_cTn.set('id', '2')
    seq_cTn.set('dur', 'indefinite')
    seq_cTn.set('nodeType', 'mainSeq')
    seq_children = etree.SubElement(seq_cTn, qn('p:childTnLst'))

    prev_lst = etree.SubElement(seq, qn('p:prevCondLst'))
    pc = etree.SubElement(prev_lst, qn('p:cond'))
    pc.set('evt', 'onPrev')
    pc.set('delay', '0')
    etree.SubElement(etree.SubElement(pc, qn('p:tgtEl')), qn('p:sldTgt'))
    next_lst = etree.SubElement(seq, qn('p:nextCondLst'))
    nc = etree.SubElement(next_lst, qn('p:cond'))
    nc.set('evt', 'onNext')
    nc.set('delay', '0')
    etree.SubElement(etree.SubElement(nc, qn('p:tgtEl')), qn('p:sldTgt'))

    ctn_id = 3
    bld_spids = []

    for group in shape_groups:
        click_par = etree.SubElement(seq_children, qn('p:par'))
        click_cTn = etree.SubElement(click_par, qn('p:cTn'))
        click_cTn.set('id', str(ctn_id)); ctn_id += 1
        click_cTn.set('fill', 'hold')
        cl = etree.SubElement(click_cTn, qn('p:stCondLst'))
        etree.SubElement(cl, qn('p:cond')).set('delay', '0')
        click_children = etree.SubElement(click_cTn, qn('p:childTnLst'))

        for shape in group:
            spid = str(shape.shape_id)
            bld_spids.append(spid)

            s_par = etree.SubElement(click_children, qn('p:par'))
            s_cTn = etree.SubElement(s_par, qn('p:cTn'))
            s_cTn.set('id', str(ctn_id)); ctn_id += 1
            s_cTn.set('fill', 'hold')
            sl = etree.SubElement(s_cTn, qn('p:stCondLst'))
            etree.SubElement(sl, qn('p:cond')).set('delay', '0')
            s_children = etree.SubElement(s_cTn, qn('p:childTnLst'))

            e_par = etree.SubElement(s_children, qn('p:par'))
            e_cTn = etree.SubElement(e_par, qn('p:cTn'))
            e_cTn.set('id', str(ctn_id)); ctn_id += 1
            e_cTn.set('presetID', '1')
            e_cTn.set('presetClass', 'entr')
            e_cTn.set('presetSubtype', '0')
            e_cTn.set('fill', 'hold')
            e_cTn.set('nodeType', 'clickEffect')
            el = etree.SubElement(e_cTn, qn('p:stCondLst'))
            etree.SubElement(el, qn('p:cond')).set('delay', '0')
            e_children = etree.SubElement(e_cTn, qn('p:childTnLst'))

            p_set = etree.SubElement(e_children, qn('p:set'))
            cBhvr = etree.SubElement(p_set, qn('p:cBhvr'))
            set_cTn = etree.SubElement(cBhvr, qn('p:cTn'))
            set_cTn.set('id', str(ctn_id)); ctn_id += 1
            set_cTn.set('dur', '1')
            set_cTn.set('fill', 'hold')
            scl = etree.SubElement(set_cTn, qn('p:stCondLst'))
            etree.SubElement(scl, qn('p:cond')).set('delay', '0')
            tgtEl = etree.SubElement(cBhvr, qn('p:tgtEl'))
            etree.SubElement(tgtEl, qn('p:spTgt')).set('spid', spid)
            anl = etree.SubElement(cBhvr, qn('p:attrNameLst'))
            etree.SubElement(anl, qn('p:attrName')).text = 'style.visibility'
            to_el = etree.SubElement(p_set, qn('p:to'))
            etree.SubElement(to_el, qn('p:strVal')).set('val', 'visible')

    bldLst = etree.SubElement(timing, qn('p:bldLst'))
    for spid in bld_spids:
        bp = etree.SubElement(bldLst, qn('p:bldP'))
        bp.set('spid', spid)
        bp.set('grpId', '0')
        bp.set('animBg', '1')


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════
# BUILD THE DECK
# ═══════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ───────────────────────────────────
# SLIDE 1 — Title
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.5), Inches(2.2), Inches(10), Inches(2.0),
   'The Next Layer\nfor Digital Twins', 54, WHITE, bold=True)
tb(s, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
   'From Simulation to Autonomous Action', 26, CREAM)
rect(s, Inches(1.5), Inches(5.4), Inches(3), Inches(0.012), fill=ACCENT)
tb(s, Inches(1.5), Inches(5.65), Inches(7), Inches(0.4),
   'Anthony Poole  ·  Software Engineer, Microsoft', 18, DIM)
tb(s, Inches(1.5), Inches(6.05), Inches(6), Inches(0.4),
   'Microelectronics US 2026  ·  Embedded Systems Stage  ·  April 22',
   14, DARKER_DIM)


# ───────────────────────────────────
# SLIDE 2 — Section Title: Digital Twins
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
   'Digital Twins', 72, WHITE, bold=True)
rect(s, Inches(1.5), Inches(4.3), Inches(3), Inches(0.012), fill=ACCENT)


# ───────────────────────────────────
# SLIDE 3 — What Is a Digital Twin? (bigger)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'The Foundation', 20, DIM)

# Large Physical System box (left)
rrect_text(s, Inches(0.8), Inches(1.8), Inches(4.8), Inches(4.0),
           'Physical\nSystem', 32, WHITE, fill=CARD_BG, line=ACCENT)

# Large arrows — arrow character on its own line, label below
tb(s, Inches(5.7), Inches(2.0), Inches(2.0), Inches(1.2),
   '→', 72, CREAM, align=PP_ALIGN.CENTER)
tb(s, Inches(5.7), Inches(2.9), Inches(2.0), Inches(0.5),
   'Sensors & Data', 16, DIM, align=PP_ALIGN.CENTER)
tb(s, Inches(5.7), Inches(3.8), Inches(2.0), Inches(1.2),
   '←', 72, CREAM, align=PP_ALIGN.CENTER)
tb(s, Inches(5.7), Inches(4.7), Inches(2.0), Inches(0.5),
   'Simulation', 16, DIM, align=PP_ALIGN.CENTER)

# Large Digital Twin box (right)
rrect_text(s, Inches(7.7), Inches(1.8), Inches(4.8), Inches(4.0),
           'Digital\nTwin', 32, WHITE, fill=CARD_BG, line=ACCENT)

tb(s, Inches(1.2), Inches(6.2), Inches(11), Inches(0.8),
   'A virtual model, continuously synced with a physical system.',
   20, DIM)


# ───────────────────────────────────
# SLIDE 4 — The Stakes (vertical stack with accent lines)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

metrics = [
    ('$22K / min', 'Automotive downtime'),
    ('$5M / hour', 'Semiconductor fab downtime'),
    ('$121B / year', 'U.S. grid outage costs'),
]
row_h = 1.7
# Center 3 rows vertically: total = 3*1.7 = 5.1, slide = 7.5, margin = (7.5-5.1)/2 = 1.2
start_y = 1.2
for i, (val, label) in enumerate(metrics):
    y = start_y + i * row_h
    # Thin accent line above each row (skip first)
    if i > 0:
        line_y = y - 0.2
        rect(s, Inches(1.5), Inches(line_y), Inches(10.3), Inches(0.02), fill=DARKER_DIM)
    # Big number — left side
    tb(s, Inches(1.5), Inches(y), Inches(7.5), Inches(1.2),
       val, 60, WHITE, bold=True)
    # Label — right aligned
    tf = tb(s, Inches(8.5), Inches(y + 0.2), Inches(3.3), Inches(0.9),
       label, 22, DIM)
    for p in tf.text_frame.paragraphs:
        p.alignment = PP_ALIGN.RIGHT


# ───────────────────────────────────
# SLIDE 5 — The Evolution (build animation)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
s5 = s  # Save reference for animations
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'The Evolution', 20, DIM)

gens = [
    ('Gen 1', 'Static', '"Here\'s what\nit looks like"'),
    ('Gen 2', 'Reactive', '"Something\nwent wrong"'),
    ('Gen 3', 'Predictive', '"Something\nwill go wrong"'),
    ('Gen 4', '?', ''),
]
card_w = Inches(2.7)
total_w = 4 * card_w.inches + 3 * 0.3
margin = (13.333 - total_w) / 2
s5_groups = []
for i, (label, name, desc) in enumerate(gens):
    x = Inches(margin + i * (card_w.inches + 0.3))
    group = []
    fill = CARD_BG if i < 3 else None
    line_c = DARKER_DIM if i < 3 else ACCENT
    card = rrect_text(s, x, Inches(1.8), Inches(2.7), Inches(1.3),
              f'{label}: {name}', 20, WHITE if i < 3 else CREAM,
              fill=fill, line=line_c, bold=True)
    group.append(card)
    if desc:
        desc_tb = tb(s, x, Inches(3.3), Inches(2.7), Inches(1.2),
           desc, 15, DIM, align=PP_ALIGN.CENTER)
        group.append(desc_tb)
    if i < 3:
        arrow_x = x + card_w
        arr = tb(s, arrow_x, Inches(2.1), Inches(0.3), Inches(1.0),
           '→', 24, DARKER_DIM, align=PP_ALIGN.CENTER)
        group.append(arr)
    s5_groups.append(group)

s5_bottom = tb(s, Inches(1.0), Inches(5.2), Inches(11), Inches(1.5),
   'Each generation gets smarter about what\'s happening.\nBut the human is still the one who decides what to do about it.',
   22, CREAM)
s5_groups.append([s5_bottom])


# ───────────────────────────────────
# SLIDE 5 — The Confession
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.5), Inches(2.5), Inches(10), Inches(3.0),
   'Digital twins can see everything.\nThey can\'t do anything\nabout what they see.',
   42, CREAM, bold=True)


# ───────────────────────────────────
# SLIDE 6 — The 2 AM Story (minimal, speaker carries it)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

tb(s, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5),
   '2:14 AM', 72, WHITE, bold=True, align=PP_ALIGN.CENTER)


# ───────────────────────────────────
# SLIDE 7 — "Agents"
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
   'Agents', 72, WHITE, bold=True)
rect(s, Inches(1.5), Inches(4.3), Inches(3), Inches(0.012), fill=ACCENT)


# ───────────────────────────────────
# SLIDE 9 — Four Capabilities (build animation)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
s9 = s  # Save reference for animations
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'What Makes an Agent', 20, DIM)

caps = [
    ('Autonomy', 'Acts within defined\nboundaries without\nprompting'),
    ('Planning', 'Breaks complex problems\ninto multi-step\nworkflows'),
    ('Tool Use', 'Connects to real\nsystems, APIs,\nand databases'),
    ('Reflection', 'Evaluates its own\nactions and adjusts\nstrategy'),
]
card_w = Inches(2.7)
total_w = 4 * card_w.inches + 3 * 0.3  # 3 gaps of 0.3"
margin = (13.333 - total_w) / 2
s9_groups = []
for i, (title, desc) in enumerate(caps):
    x = Inches(margin + i * (card_w.inches + 0.3))
    group = []
    card = rrect_text(s, x, Inches(2.0), card_w, Inches(1.2),
              title, 22, WHITE, fill=CARD_BG, line=ACCENT, bold=True)
    group.append(card)
    desc_tb = tb(s, x + Inches(0.15), Inches(3.5), Inches(2.5), Inches(1.8),
       desc, 15, DIM, align=PP_ALIGN.CENTER)
    group.append(desc_tb)
    s9_groups.append(group)


# ───────────────────────────────────
# SLIDE 9 — The Convergence (fixed Venn)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

# Two circles — outline only (no fill) so intersection outlines are both visible
# Left circle: DT
oval(s, Inches(1.8), Inches(1.2), Inches(5.5), Inches(5.5),
     fill=None, line=ACCENT, lw=Pt(1.5))
tb(s, Inches(2.3), Inches(3.2), Inches(3.5), Inches(1.0),
   'Digital\nTwins', 30, WHITE, bold=True)
tb(s, Inches(2.3), Inches(4.5), Inches(3.2), Inches(1.0),
   'The data\nThe models\nThe simulation', 16, DIM)

# Right circle: Agentic AI
oval(s, Inches(6.0), Inches(1.2), Inches(5.5), Inches(5.5),
     fill=None, line=ACCENT, lw=Pt(1.5))
tb(s, Inches(7.8), Inches(3.2), Inches(3.5), Inches(1.0),
   'Agentic\nAI', 30, WHITE, bold=True)
tb(s, Inches(7.8), Inches(4.5), Inches(3.2), Inches(1.0),
   'The reasoning\nThe planning\nThe action', 16, DIM)

# Clear ? in the visible intersection zone
tb(s, Inches(5.5), Inches(3.3), Inches(2.3), Inches(1.5),
   '?', 56, CREAM, bold=True, align=PP_ALIGN.CENTER)


# ───────────────────────────────────
# SLIDE 10 — The Reveal (transition)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

tb(s, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5),
   'What if the twin could think?', 56, WHITE, bold=True, align=PP_ALIGN.CENTER)


# ───────────────────────────────────
# SLIDE 11 — The Loop (control-loop style)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'The Agentic Control Loop', 20, DIM)

# Horizontal control loop — more familiar to engineers
# Perceive → Reason → Plan → Act → Reflect (with feedback arrow back)
steps = ['Perceive', 'Reason', 'Plan', 'Act', 'Reflect']
box_w = Inches(2.0)
box_h = Inches(1.2)
start_x = Inches(0.8)
y_main = Inches(3.0)
gap = Inches(0.5)

for i, step in enumerate(steps):
    x = start_x + i * (box_w + gap)
    rrect_text(s, x, y_main, box_w, box_h,
              step, 20, WHITE, fill=CARD_BG, line=ACCENT)
    # Arrow to next (except last) — vertically centered on box
    if i < 4:
        ax = x + box_w + Inches(0.05)
        arrow_h = Inches(0.6)
        arrow_y = y_main + (box_h - arrow_h) / 2
        tb(s, ax, arrow_y, gap - Inches(0.1), arrow_h,
           '→', 22, ACCENT, align=PP_ALIGN.CENTER)

# Feedback arrow: Reflect back to Perceive (drawn below the boxes)
fb_y = y_main + box_h + Inches(0.3)
last_x = start_x + 4 * (box_w + gap) + box_w / 2
first_x = start_x + box_w / 2

# Bottom feedback line
rect(s, first_x, fb_y, last_x - first_x, Inches(0.015), fill=ACCENT)
# Right vertical (down from Reflect)
rect(s, last_x - Inches(0.008), y_main + box_h, Inches(0.015),
     fb_y - y_main - box_h + Inches(0.015), fill=ACCENT)
# Left vertical (up to Perceive)
rect(s, first_x - Inches(0.008), y_main + box_h, Inches(0.015),
     fb_y - y_main - box_h + Inches(0.015), fill=ACCENT)
# Arrowhead pointing up at Perceive box
tb(s, first_x - Inches(0.25), y_main + box_h - Inches(0.15), Inches(0.5), Inches(0.4),
   '▲', 14, ACCENT, align=PP_ALIGN.CENTER)

# Label on feedback
tb(s, Inches(4.5), fb_y + Inches(0.1), Inches(4.5), Inches(0.5),
   'Continuous feedback', 14, DIM, align=PP_ALIGN.CENTER)

# Subtitle
tb(s, Inches(1.2), Inches(5.5), Inches(11), Inches(1.0),
   'Perceive the environment. Reason about causes. Plan a response.\nAct through connected systems. Reflect on outcomes.',
   18, CREAM)


# ───────────────────────────────────
# SLIDE 12 — Architecture (simplified three layers)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

layers = [
    ('The Mind',       'Cognitive Layer',  ACCENT,     CREAM),
    ('The Mirror',     'Digital Twin',     DARKER_DIM, WHITE),
    ('The Real Thing', 'Physical System',  DARKER_DIM, WHITE),
]

for i, (label, name, line_c, text_c) in enumerate(layers):
    y = Inches(1.3 + i * 1.7)
    h = Inches(1.3)
    rrect_text(s, Inches(3.2), y, Inches(8.5), h,
              name, 28, text_c, fill=CARD_BG, line=line_c)
    tb(s, Inches(0.6), y + Inches(0.3), Inches(2.4), Inches(0.7),
       label, 22, CREAM if i == 0 else DIM, bold=i == 0,
       align=PP_ALIGN.RIGHT)


# ───────────────────────────────────
# SLIDE 13 — The Evidence
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'The Evidence', 20, DIM)

evidence = [
    ('IBM — AAAI 2025',
     'First major academic paper on agentic digital twins.\n'
     'LLM agents for shipping fleet management —\n'
     'autonomously plan, execute, reflect.'),
    ('XMPro MAGS',
     'Production multi-agent system for manufacturing.\n'
     'Potential 30–40% reduction in unplanned downtime.'),
    ('Nature Computational Science (2026)',
     'Mapped Digital Twin evolution from reactive to agentic,\n'
     'naming six capability levels: Standalone · Descriptive ·\n'
     'Diagnostic · Predictive · Prescriptive · Autonomous'),
]
for i, (title, desc) in enumerate(evidence):
    y = Inches(1.4 + i * 1.9)
    rect(s, Inches(1.5), y, Inches(0.08), Inches(1.4), fill=ACCENT)
    tb(s, Inches(2.0), y, Inches(9), Inches(0.5), title, 22, WHITE, bold=True)
    tb(s, Inches(2.0), y + Inches(0.5), Inches(9), Inches(1.0), desc, 16, DIM)


# ───────────────────────────────────
# SLIDE 14 — Multi-Agent (simplified)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(8), Inches(0.5),
   'Multi-Agent Systems', 20, DIM)

# Three agents — simpler, just cards with role name
agents = [
    ('Monitor', 'Detects anomalies'),
    ('Analyst', 'Finds root causes'),
    ('Planner', 'Schedules responses'),
]
card_w = Inches(3.5)
total_w = 3 * card_w.inches + 2 * 0.4  # 2 gaps of 0.4"
margin = (13.333 - total_w) / 2
for i, (name, desc) in enumerate(agents):
    x = Inches(margin + i * (card_w.inches + 0.4))
    rrect_text(s, x, Inches(2.2), card_w, Inches(2.0),
              name, 28, WHITE, fill=CARD_BG, line=ACCENT)
    tb(s, x, Inches(4.4), card_w, Inches(0.6),
       desc, 17, DIM, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow_h = Inches(0.8)
        arrow_y = Inches(2.2) + (Inches(2.0) - arrow_h) / 2
        tb(s, x + card_w, arrow_y, Inches(0.4), arrow_h,
           '→', 28, ACCENT, align=PP_ALIGN.CENTER)

tb(s, Inches(1.2), Inches(5.6), Inches(11), Inches(1.0),
   'Specialized agents. Collaborate and check each other\'s work.\n'
   'The agent that proposes is never the one that approves.',
   18, CREAM)


# ───────────────────────────────────
# SLIDE 15 — Three Worlds, One Pattern
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(8), Inches(0.5),
   'One Pattern, Every System', 20, DIM)

domains = [
    ('Edge AI Device', 'Model drift\n→ Diagnose\n→ Retrain'),
    ('Robotic Arm', 'Bearing wear\n→ Predict\n→ Schedule'),
    ('ADAS Sensor Suite', 'Sensor fusion\n→ Safety\n→ Adapt'),
]
card_w = Inches(3.6)
total_w = 3 * card_w.inches + 2 * 0.35
margin = (13.333 - total_w) / 2
for i, (name, flow) in enumerate(domains):
    x = Inches(margin + i * (card_w.inches + 0.35))
    rrect_text(s, x, Inches(1.8), card_w, Inches(1.3),
              name, 24, WHITE, fill=CARD_BG, line=ACCENT)
    tb(s, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(2.0),
       flow, 18, DIM, align=PP_ALIGN.LEFT)

tb(s, Inches(1.0), Inches(5.8), Inches(11), Inches(1.0),
   'Same loop. Same architecture. Different worlds.\nThe pattern is universal.',
   22, CREAM, bold=True)


# ───────────────────────────────────
# SLIDE 16 — Challenges (fixed overlap)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)
tb(s, Inches(1.2), Inches(0.5), Inches(5), Inches(0.5),
   'Let\'s Be Honest', 20, DIM)

challenges = [
    ('Accountability', 'When the agent decides wrong, who owns it?'),
    ('Explainability', 'Can you trace why it made that call?'),
    ('Security', 'New attack surfaces we haven\'t fully mapped.'),
    ('Governance', 'Our frameworks weren\'t built for autonomous agents.'),
]
for i, (title, desc) in enumerate(challenges):
    y = Inches(1.4 + i * 1.3)
    rect(s, Inches(1.5), y + Inches(0.05), Inches(0.08), Inches(0.8), fill=ACCENT)
    tb(s, Inches(2.0), y, Inches(4), Inches(0.5), title, 22, WHITE, bold=True)
    tb(s, Inches(2.0), y + Inches(0.45), Inches(9), Inches(0.5), desc, 17, DIM)


# ───────────────────────────────────
# SLIDE 17 — Demo CTA
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

tb(s, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
   'Talk to an agentic digital twin yourself.', 36, WHITE, bold=True)

# QR code — demo
qr_demo = os.path.join(SCRIPT_DIR, 'qr-demo.png')
if os.path.exists(qr_demo):
    s.shapes.add_picture(qr_demo, Inches(4.7), Inches(2.5), Inches(3.9), Inches(3.5))
else:
    rrect_text(s, Inches(4.5), Inches(2.5), Inches(4.3), Inches(3.5),
              '[ QR Code ]', 24, DIM, fill=CARD_BG, line=ACCENT, bold=False)

# Three domain labels
for i, name in enumerate(['Edge AI', 'Robotic Arm', 'ADAS']):
    x = Inches(2.5 + i * 3.2)
    rrect_text(s, x, Inches(6.3), Inches(2.5), Inches(0.6),
              name, 14, CREAM, fill=CARD_BG, line=DARKER_DIM, bold=False)

tb(s, Inches(0), Inches(6.95), Inches(13.333), Inches(0.4),
   'tonypdemo.it.com', 18, DIM, align=PP_ALIGN.CENTER)


# ───────────────────────────────────
# SLIDE 19 — Questions + Contact
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

tb(s, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
   'Questions?', 52, WHITE, bold=True)

rect(s, Inches(1.5), Inches(2.8), Inches(2.5), Inches(0.012), fill=ACCENT)

tb(s, Inches(1.5), Inches(3.3), Inches(6), Inches(1.0),
   'Anthony Poole', 36, WHITE, bold=True)
tb(s, Inches(1.5), Inches(4.1), Inches(6), Inches(0.6),
   'Software Engineer, Microsoft', 22, CREAM)
tb(s, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
   'anthony50102.github.io', 20, DIM)
tb(s, Inches(1.5), Inches(5.3), Inches(10), Inches(0.5),
   'linkedin.com/in/anthony-poole-079548206', 20, DIM)

# Personal site QR code (right side)
qr_personal = os.path.join(SCRIPT_DIR, 'qr-personal.png')
if os.path.exists(qr_personal):
    s.shapes.add_picture(qr_personal, Inches(9.2), Inches(2.8), Inches(3.2), Inches(3.2))
    tb(s, Inches(9.0), Inches(6.05), Inches(3.6), Inches(0.4),
       'anthony50102.github.io', 16, DIM, align=PP_ALIGN.CENTER)
else:
    rrect_text(s, Inches(9.5), Inches(3.3), Inches(2.8), Inches(2.5),
              '[ QR Code ]\nanthonypoole.dev', 14, DIM, fill=CARD_BG, line=DARKER_DIM, bold=False)


# ───────────────────────────────────
# SLIDE 20 — Close (final takeaway)
# ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(s)

tb(s, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
   'Digital twins give us the body.\nIt\'s time to consider the mind.',
   32, CREAM, italic=True)

rect(s, Inches(1.5), Inches(5.3), Inches(2.5), Inches(0.012), fill=ACCENT)

tb(s, Inches(1.5), Inches(5.6), Inches(6), Inches(0.5),
   'Thank you.', 22, CREAM)


# ═══════════════════════════════════════
# POST-BUILD: Animations
# ═══════════════════════════════════════
add_appear_animations(s5, s5_groups)
add_appear_animations(s9, s9_groups)


# ═══════════════════════════════════════
# POST-BUILD: Speaker Notes
# ═══════════════════════════════════════
speaker_notes = {
    0: (  # Slide 1 — Title
        "Hey everyone. I'm Anthony Poole, software engineer at Microsoft — "
        "I build AI systems. But I actually started a bit closer to y'all's world "
        "than you'd think — I interned at AMD, did research in photonics, wrote my "
        "thesis on digital twins, and I worked on embedded systems for autonomous "
        "vehicles. But then I crossed over to the big tech AI side. And over the "
        "last year I've watched something reshape how my team thinks about every "
        "system we touch. Now when I look back at what's happening with digital "
        "twins, I realize the same shift is heading this way."
    ),
    1: (  # Slide 2 — Section: Digital Twins
        "[Section divider — click through quickly. This is a signpost, not a beat.]"
    ),
    2: (  # Slide 3 — What Is a Digital Twin?
        "So let's make sure we're starting from the same place. A digital twin — "
        "at its core — is a virtual model that stays continuously synced with a "
        "physical system. Data flows in from sensors. Simulations flow back. "
        "\"Most of you know this better than I do. I'm not here to explain digital "
        "twins to you. I'm here to talk about what comes next.\""
    ),
    3: (  # Slide 4 — The Stakes
        "In automotive manufacturing, unplanned downtime costs $22,000 a minute. "
        "In semiconductor fabs it can hit $5 million an hour. And at the electric "
        "grid level, U.S. power outages cost $121 billion last year alone. "
        "Here's the real problem: these systems are generating more data every "
        "minute than teams can act on in a week."
    ),
    4: (  # Slide 5 — The Evolution
        "[BUILD ANIMATION — click through generations one at a time.]\n\n"
        "Digital twins have been evolving. Gen 1 was static — a 3D model, a "
        "schematic. Gen 2 was reactive — something went wrong, here's an alert. "
        "Gen 3 — where a lot of you are or are heading — is predictive.\n\n"
        "But notice — every generation gets smarter about what's happening. "
        "None of them do anything about it."
    ),
    5: (  # Slide 6 — The Confession
        "[Let the text land for a beat.]\n\n"
        "Digital twins can see everything. They can't do anything about what "
        "they see. That's not a failure — that's a ceiling. And I think we're "
        "about to break through it."
    ),
    6: (  # Slide 7 — 2:14 AM
        "[Slow. Conversational. Like you're telling a friend.]\n\n"
        "You're a process engineer at a fab. It's 2:14 AM. Something shifts — "
        "temperature drift, vibration spike. Your monitoring system catches it "
        "instantly. You? You're asleep. By the time you wake up, understand "
        "what happened, track down the part — five, six hours have passed.\n\n"
        "The system detected it in seconds. It took the better part of a shift. "
        "Detection and action are two completely different problems."
    ),
    7: (  # Slide 8 — Section: Agents
        "That was digital twins. Now let's talk about agents.\n\n"
        "I know — you've heard the word. It's been on keynote slides and "
        "LinkedIn posts for a year. But I've spent the last year building "
        "these systems, and they work. They've changed how my team solves "
        "problems and makes decisions."
    ),
    8: (  # Slide 9 — Four Capabilities
        "[BUILD ANIMATION — click through capabilities one at a time.]\n\n"
        "Autonomy — it operates within boundaries you define, without waiting. "
        "Planning — it breaks complex problems into multi-step workflows. "
        "Tool use — it connects to real systems, APIs, databases. "
        "Reflection — it evaluates what it did and adjusts its approach.\n\n"
        "That last one is the one people underestimate."
    ),
    9: (  # Slide 10 — Convergence Venn
        "On one side, digital twins — the data, the models, the simulation. "
        "On the other, agentic AI — the reasoning, the planning, the action. "
        "If you're looking at this thinking 'why aren't these connected yet?' "
        "— that's exactly the right question."
    ),
    10: (  # Slide 11 — The Reveal
        "[Let the slide breathe. A beat of silence before you speak.]\n\n"
        "What if the twin could think? A digital twin that doesn't just "
        "mirror a system — it reasons about it. It plans. It acts. Not "
        "replacing the engineer. Giving the engineer a partner that never "
        "sleeps and never loses context."
    ),
    11: (  # Slide 12 — The Agentic Control Loop
        "It's a control loop — something every engineer in this room already "
        "thinks in. Perceive — ingest data. Reason — identify what's happening "
        "and why. Plan — build a response. Act — execute through connected "
        "systems. Reflect — evaluate the outcome.\n\n"
        "Remember that engineer at 2 AM? An agentic twin wouldn't have just "
        "detected the drift. It would have pulled historical patterns, "
        "identified the failing element, and scheduled the repair. Minutes, "
        "not hours."
    ),
    12: (  # Slide 13 — Architecture
        "Architecturally, it's simpler than you'd think. Physical system at "
        "the bottom. Digital twin in the middle. And on top, a new cognitive "
        "layer — the reasoning and decision-making.\n\n"
        "Two of those three layers? You already have them. The whole idea is "
        "one new layer on top of what you've already invested in."
    ),
    13: (  # Slide 14 — The Evidence
        "IBM published a paper at AAAI 2025 on agentic digital twins for "
        "shipping fleet management. XMPro has a multi-agent system in "
        "production — reporting 30–40% reductions in unplanned downtime. "
        "And Nature Computational Science published a formal evolutionary "
        "framework this year.\n\n"
        "This went from idea to research to production faster than most "
        "people realize."
    ),
    14: (  # Slide 15 — Multi-Agent
        "The architecture that's emerging isn't one giant AI doing everything. "
        "It's specialized agents working together. A monitor that detects. "
        "An analyst that finds root causes. A planner that schedules the "
        "response. They collaborate — and critically — they check each "
        "other's work. The agent that proposes is never the one that approves.\n\n"
        "If that sounds like how a good engineering team works — it is."
    ),
    15: (  # Slide 16 — Three Worlds, One Pattern
        "This isn't a solution for one industry. The same pattern applies "
        "everywhere. Edge AI device drifts? The twin diagnoses and retrains. "
        "Robotic arm shows bearing wear? The twin predicts and schedules "
        "maintenance. ADAS sensor gets conflicting data? The twin runs safety "
        "checks and adapts.\n\n"
        "Different worlds. Same idea. That's how you know it's real."
    ),
    16: (  # Slide 17 — Challenges
        "When an agent makes a bad decision, who owns it? Can you trace why "
        "it made that call? These systems create attack surfaces we haven't "
        "fully mapped. And governance frameworks weren't built for tools "
        "that act on their own.\n\n"
        "These are solvable problems. But they're real."
    ),
    17: (  # Slide 18 — Demo CTA
        "I built something. Scan the QR code — you can talk to an agentic "
        "digital twin right now. Three of them: edge AI, robotic arm, ADAS. "
        "Ask what's wrong. Inject issues. Watch it reason.\n\n"
        "[Give them 10–15 seconds to scan.]\n\n"
        "The URL is at the bottom if you'd rather type. It'll be live after "
        "the talk too — no rush."
    ),
    18: (  # Slide 19 — Questions
        "Digital twins already have the data and models. Agentic AI gives "
        "them the ability to act on it. That's it. That's the whole idea. "
        "One new layer.\n\n"
        "Thank you. Come find me or scan the QR for my site.\n\n"
        "[This slide stays up during Q&A — contact info and QR visible.]"
    ),
    19: (  # Slide 20 — Close
        "[Click to this after Q&A wraps, as your final word.]\n\n"
        "Thanks everyone."
    ),
}

for idx, notes_text in speaker_notes.items():
    if idx < len(prs.slides):
        add_notes(prs.slides[idx], notes_text)


# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
out = '/Users/anthonypoole/Repositories/MicroUSPresentation/presentation/presentation-v4-gradient-noir.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Total slides: {len(prs.slides)}')
